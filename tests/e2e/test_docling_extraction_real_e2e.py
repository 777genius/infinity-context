import asyncio
import json
from io import BytesIO
from pathlib import Path

import pytest
from infinity_context_adapters.extraction.docling_engine import DoclingDocumentExtractionEngine
from infinity_context_core.ports.extraction import ExtractionLimits, ExtractionRequest
from pypdf import PdfReader
from resource_guards import apply_stable_ml_env, skip_if_low_temp_space

pytest.importorskip("docling")
pytest.importorskip("docx")
pytest.importorskip("openpyxl")
pytest.importorskip("PIL")
pytest.importorskip("pptx")


def test_sample_pdf_fixture_is_strictly_parseable() -> None:
    expected_text = "PDF Memory Scope Evidence"

    reader = PdfReader(BytesIO(_sample_pdf_bytes(expected_text)), strict=True)

    assert len(reader.pages) == 1
    assert reader.pages[0].extract_text().strip() == expected_text


def test_real_docling_extracts_document_evidence_and_artifacts(tmp_path: Path) -> None:
    apply_stable_ml_env(tmp_path / "tmp")
    skip_if_low_temp_space(tmp_path, label="Real Docling extraction e2e")
    cases = (
        _case(
            filename="docling-e2e.pdf",
            content_type="application/pdf",
            content=_sample_pdf_bytes("PDF Memory Scope Evidence"),
            expected_text="PDF Memory Scope Evidence",
            expected_page_refs=True,
            expected_bbox_refs=True,
        ),
        _case(
            filename="docling-e2e.docx",
            content_type=(
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            ),
            content=_sample_docx_bytes(),
            expected_text="Alex approved frontend quick capture",
            expected_table_html=True,
        ),
        _case(
            filename="docling-e2e.pptx",
            content_type=(
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ),
            content=_sample_pptx_bytes(),
            expected_text="PPTX Image Evidence",
            expected_page_refs=True,
            expected_bbox_refs=True,
            expected_image_refs=True,
        ),
        _case(
            filename="docling-e2e.xlsx",
            content_type=("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
            content=_sample_xlsx_bytes(),
            expected_text="Keep capture evidence linked",
            expected_table_html=True,
            expected_page_refs=True,
            expected_bbox_refs=True,
        ),
    )

    for case in cases:
        result = asyncio.run(DoclingDocumentExtractionEngine().extract(case["request"]))

        assert result.status == "succeeded", case["filename"]
        assert result.parser_name == "docling_document"
        assert case["expected_text"] in (result.markdown or "")
        assert result.elements
        assert result.technical_metadata["docling_element_strategy"] == "docling_items"
        assert "normalized_json" in _artifact_types(result.artifacts)
        normalized = next(
            item for item in result.artifacts if item.artifact_type == "normalized_json"
        )
        normalized_payload = json.loads(normalized.content.decode("utf-8"))
        assert normalized_payload.get("schema_name") or normalized_payload.get("body")

        if case["expected_table_html"]:
            table_html = [item for item in result.artifacts if item.artifact_type == "table_html"]
            assert table_html, case["filename"]
            assert b"<table" in table_html[0].content
            assert result.technical_metadata["docling_table_count"] >= 1

        if case["expected_image_refs"]:
            assert any(item.kind == "image" for item in result.elements), case["filename"]
            assert result.technical_metadata["docling_image_count"] >= 1

        if case["expected_page_refs"]:
            assert result.technical_metadata["docling_page_ref_count"] >= 1
            assert any(item.page_number is not None for item in result.elements)

        if case["expected_bbox_refs"]:
            assert result.technical_metadata["docling_bbox_ref_count"] >= 1
            assert any(item.bbox is not None for item in result.elements)


def _case(
    *,
    filename: str,
    content_type: str,
    content: bytes,
    expected_text: str,
    expected_table_html: bool = False,
    expected_image_refs: bool = False,
    expected_page_refs: bool = False,
    expected_bbox_refs: bool = False,
) -> dict[str, object]:
    return {
        "filename": filename,
        "expected_text": expected_text,
        "expected_table_html": expected_table_html,
        "expected_image_refs": expected_image_refs,
        "expected_page_refs": expected_page_refs,
        "expected_bbox_refs": expected_bbox_refs,
        "request": ExtractionRequest(
            job_id="job-docling-e2e",
            asset_id="asset-docling-e2e",
            filename=filename,
            declared_content_type=content_type,
            detected_content_type=content_type,
            byte_size=len(content),
            sha256_hex="0" * 64,
            content=content,
            parser_profile="standard_docling",
            limits=ExtractionLimits(
                max_bytes=10_000_000,
                max_pages=20,
                max_output_chars=20_000,
                max_tables=5,
            ),
        ),
    }


def _artifact_types(artifacts: object) -> set[str]:
    return {item.artifact_type for item in artifacts}


def _sample_pdf_bytes(text: str) -> bytes:
    escaped_text = (
        text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)").encode("latin-1")
    )
    stream = b"BT\n/F1 18 Tf\n72 720 Td\n(" + escaped_text + b") Tj\nET\n"
    objects = (
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length "
        + str(len(stream)).encode("ascii")
        + b" >>\nstream\n"
        + stream
        + b"endstream",
    )
    document = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets: list[int] = []
    for number, content in enumerate(objects, start=1):
        offsets.append(len(document))
        document.extend(f"{number} 0 obj\n".encode("ascii"))
        document.extend(content)
        document.extend(b"\nendobj\n")

    xref_offset = len(document)
    document.extend(b"xref\n0 6\n0000000000 65535 f \n")
    for offset in offsets:
        document.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    document.extend(b"trailer\n<< /Size 6 /Root 1 0 R >>\nstartxref\n")
    document.extend(f"{xref_offset}\n%%EOF\n".encode("ascii"))
    return bytes(document)


def _sample_docx_bytes() -> bytes:
    from docx import Document

    document = Document()
    document.add_heading("DOCX Memory Scope Evidence", 1)
    document.add_paragraph("Alex approved frontend quick capture.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Scope"
    table.cell(0, 1).text = "Thread"
    table.cell(1, 0).text = "frontend"
    table.cell(1, 1).text = "alex-call"
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _sample_pptx_bytes() -> bytes:
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "PPTX Image Evidence"
    image = Image.new("RGB", (60, 30), color=(255, 255, 255))
    image_buffer = BytesIO()
    image.save(image_buffer, format="PNG")
    image_buffer.seek(0)
    slide.shapes.add_picture(image_buffer, Inches(1), Inches(2), width=Inches(1))
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _sample_xlsx_bytes() -> bytes:
    from openpyxl import Workbook

    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Memory"
    worksheet.append(["Scope", "Decision"])
    worksheet.append(["frontend", "Keep capture evidence linked"])
    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()
